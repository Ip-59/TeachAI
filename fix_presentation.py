import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

TEMPLATE_PATH = "Шаблон презентации.pptx"
OUTPUT_PATH = "TeachAI_Stage4.pptx"

def set_shape_text(shape, text, font_size=None, bold=None, color=None, alignment=None):
    if not shape.has_text_frame: return
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs: run.text = ""
    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.text = text
    else:
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()
        run.text = text
    if font_size is not None: run.font.size = Pt(font_size)
    if bold is not None: run.font.bold = bold
    if color is not None: run.font.color.rgb = color
    if alignment is not None: tf.paragraphs[0].alignment = alignment

def set_multiline_text(shape, lines, font_size=None, color=None, bold=None, line_spacing=None):
    if not shape.has_text_frame: return
    tf = shape.text_frame
    existing_font_size = existing_color = existing_bold = existing_font_name = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        r = tf.paragraphs[0].runs[0]
        existing_font_size = r.font.size
        try: existing_color = r.font.color.rgb if r.font.color and r.font.color.type is not None else None
        except AttributeError: existing_color = None
        existing_bold = r.font.bold
        existing_font_name = r.font.name
    for i in range(len(tf.paragraphs) - 1, -1, -1):
        if i > 0:
            p = tf.paragraphs[i]
            p._p.getparent().remove(p._p)
    first_para = tf.paragraphs[0]
    for run in first_para.runs: run.text = ""
    for i, line in enumerate(lines):
        if i == 0:
            p = first_para
            if p.runs: run = p.runs[0]
            else: run = p.add_run()
            run.text = line
        else:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
        fs = font_size if font_size is not None else existing_font_size
        if fs is not None: run.font.size = Pt(fs) if isinstance(fs, (int, float)) else fs
        c = color if color is not None else existing_color
        if c is not None: run.font.color.rgb = c
        b = bold if bold is not None else existing_bold
        if b is not None: run.font.bold = b
        if existing_font_name: run.font.name = existing_font_name
        if line_spacing is not None: p.line_spacing = Pt(line_spacing)

def add_clean_textbox(slide, text_lines, left, top, width, height, font_size=18, bold_first=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold) in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = is_bold
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

def clear_picture_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or str(shape.shape_type) == "PLACEHOLDER (14)":
            sp = shape._element
            try: sp.getparent().remove(sp)
            except AttributeError: pass
            
def center_expand_text_shapes(slide):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text.strip():
            # If it looks like content
            if "Текст" not in shape.text and shape.top > Inches(2):
                shape.left = Inches(1.0)
                shape.width = Inches(8.0)

def create_presentation():
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)
    
    # 1. Заголовок
    slide = slides[0]
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Название" in shape.text: set_shape_text(shape, "TeachAI — Интеллектуальная система\nперсонализированного обучения", font_size=24, bold=True)
            elif "ФИО" in shape.text: set_shape_text(shape, "Игорь Пучкин", font_size=16)
            elif "Дата" in shape.text: set_shape_text(shape, "Февраль 2026", font_size=14)
            elif "Поток" in shape.text: set_multiline_text(shape, ["\n", "AI/ML — разработчик"], font_size=14)
            
    # 2. Постановка задачи
    slide = slides[1]
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Постановка" in shape.text: set_shape_text(shape, "1. Тема и описание задачи", bold=True)
            elif shape.text.startswith("Текст"):
                set_multiline_text(shape, [
                    "Проблема:\nОтсутствие доступных инструментов для персонализированного обучения программированию с адаптацией под студента.\n",
                    "Решение:\nTeachAI — интеллектуальная образовательная система, использующая GPT модель через OpenAI API для автоматической генерации персонализированных учебных материалов.\n",
                    "Тип задачи:\nAI-система прикладного уровня, использующая inference предобученной LLM для генерации контента.\n",
                    "Подход:\nPrompt Engineering вместо fine-tuning."
                ], font_size=16)
                shape.width = Inches(11.5)
                
    # 3. Цель
    slide = slides[2]
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Цель" in shape.text: set_shape_text(shape, "Цель проекта", bold=True)
            elif shape.text.startswith("Текст"):
                set_multiline_text(shape, ["Автоматизация создания персонализированного", "образовательного контента с использованием", "современных языковых моделей (LLM)", "для обеспечения эффективного", "и гибкого обучения в интерактивной среде", "Jupyter Notebook."], font_size=20)
                shape.width = Inches(11.5)
                
    # 4. Задачи
    slide = slides[3]
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Задачи" in shape.text: set_shape_text(shape, "Задачи проекта", bold=True)
            elif shape.text.strip().startswith("1"):
                set_multiline_text(shape, [
                    "1. Генерация учебных планов, уроков, примеров и тестов через OpenAI GPT API", "",
                    "2. Персонализация контента под стиль общения и параметры обучения студента", "",
                    "3. Многоуровневая валидация и автокоррекция сгенерированного контента", "",
                    "4. Интерактивный интерфейс на ipywidgets с навигацией, тестированием и Q&A", "",
                    "5. Отслеживание прогресса обучения и сохранение состояния"
                ], font_size=16)
                shape.width = Inches(11.5)
                
    # 5. Обучающая база
    slide = slides[4]
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "база" in shape.text.lower(): set_shape_text(shape, "2. База данных / Данные", bold=True)
            elif "Объем" in shape.text or "информация" in shape.text.lower():
                set_multiline_text(shape, [
                    "Сбор базы данных не требуется.\nПроект использует готовую предобученную модель.\n",
                    "Источник данных:\nOpenAI GPT-4 API — генерация контента в реальном времени\n",
                    "Вместо датасета:\nPrompt Engineering — управление качеством через промпты и самопроверку (Self-Correction)\n",
                    "Обоснование:\n• Нет задач классификации/детекции, требующих дообучения\n• Контент генерируется динамически, что даёт гибкость и актуальность знаний модели."
                ], font_size=15)
                shape.width = Inches(11.5)
                
    # 6. Архитектура
    slide = slides[5]
    clear_picture_placeholders(slide)
    for shape in slide.shapes:
        if hasattr(shape, 'text') and "картинк" in shape.text.lower(): set_shape_text(shape, "Архитектура системы TeachAI", bold=True)
    add_clean_textbox(slide, [
        ("Многослойная фасадная архитектура системы:", True),
        ("", False),
        ("1. Базовый уровень инфраструктуры:", True),
        ("   • StateManager — управление персистентным состоянием (прогресс)", False),
        ("   • ConfigManager — безопасное управление ключами и окружением", False),
        ("", False),
        ("2. Оркестрация:", True),
        ("   • TeachAIEngine — главный контроллер и инициализатор жизненного цикла", False),
        ("", False),
        ("3. Ядро генерации контента (LLM интеграция):", True),
        ("   • ContentGenerator — Фасад для взаимодействия с OpenAI", False),
        ("   • Набор генераторов: CoursePlanGenerator, LessonGenerator, AssessmentGenerator", False),
        ("", False),
        ("4. Модуль представления (UI):", True),
        ("   • UserInterface на базе ipywidgets для работы прямо в Jupyter-ноутбуке", False)
    ], Inches(1), Inches(1.5), Inches(11), Inches(5), font_size=16)

    # 7. Prompt Engineering
    slide = slides[6]
    clear_picture_placeholders(slide)
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "картинк" in shape.text.lower(): set_shape_text(shape, "Prompt Engineering — методология", bold=True)
            elif shape.text.startswith("Текст"):
                set_multiline_text(shape, [
                    "Система опирается на строгие систематизированные промпты, разделенные по зонам ответственности:",
                    "",
                    "✅ CoursePlanGenerator: формирует JSON структуру курса",
                    "✅ LessonGenerator: генерирует тело урока с разметкой Markdown",
                    "✅ ExamplesGenerator: конструирует изолированные фрагменты Python-кода",
                    "✅ AssessmentGenerator: разрабатывает тестовые задания",
                    "✅ RelevanceChecker: фильтрует нерелевантные вопросы пользователя",
                    "",
                    "Каждый генератор использует свой усиленный промпт с правилами форматирования, примерами ошибок (""Few-Shot"") и эмоджи-маркерами для акцентирования внимания LLM на важных ограничениях."
                ], font_size=18, line_spacing=24)
                shape.left = Inches(1)
                shape.width = Inches(11)

    # 8. Параметризация данных
    slide = slides[7]
    clear_picture_placeholders(slide)
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "картинк" in shape.text.lower(): set_shape_text(shape, "Параметризация данных (Контекст)", bold=True)
            elif "Текст" in shape.text:
                set_multiline_text(shape, [
                    "Параметры инференса LLM API:",
                    "• model: gpt-3.5-turbo / gpt-4",
                    "• temperature: 0.7 (для творческой генерации) / 0.2 (для строгих JSON или проверок)",
                    "• max_tokens: 3500",
                    "• response_format: json_object (гарантия получения словарей Python)",
                    "",
                    "Параметры профиля пользователя (внедряются в промпт):",
                    "• Стиль общения: formal (академичный), friendly (дружелюбный), casual (разговорный)",
                    "• Длительность планируемого урока (влияет на объем генерации)",
                    "",
                    "Удержание состояния:",
                    "• Локальный JSON файл (state.json) действует как база данных для хранения профилей, баллов и истории."
                ], font_size=16, line_spacing=22)
                shape.left = Inches(1)
                shape.width = Inches(11)

    # 9. Жизненный цикл урока
    slide = slides[8]
    clear_picture_placeholders(slide)
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "картинк" in shape.text.lower(): set_shape_text(shape, "Жизненный цикл интерактивного урока", bold=True)
            elif "Текст" in shape.text:
                set_multiline_text(shape, [
                    "Обучение в TeachAI — это непрерывный цикл:",
                    "",
                    "1. Инициализация урока: Сбор метаданных из плана курса.",
                    "2. Динамическая генерация: Вызов API OpenAI и рендеринг UI.",
                    "3. Взаимодействие (Chat): Студент может потребовать дополнительные примеры или задать уточняющий вопрос (через RelevanceChecker).",
                    "4. Контроль знаний (Assessment): Генерация 5 тестовых вопросов с принудительной рандомизацией.",
                    "5. Верификация: Оценка ответов. Если пройден порог, студент допускается к следующему модулю.",
                    "6. Сохранение прогресса: Запись результатов студента в профиль."
                ], font_size=18, line_spacing=26)
                shape.left = Inches(1)
                shape.width = Inches(11)

    # 10. Интерфейс UX
    slide = slides[9]
    clear_picture_placeholders(slide)
    for shape in slide.shapes:
        if hasattr(shape, 'text') and ("подпись" in shape.text.lower() or "точка" in shape.text.lower()):
            sp = shape._element
            sp.getparent().remove(sp)
    
    title_shape = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape.top < Inches(1.5): title_shape = shape; break
    if title_shape: set_shape_text(title_shape, "Пользовательский интерфейс (UX в Jupyter)", bold=True)
    
    add_clean_textbox(slide, [
        ("Интеграция с ipywidgets:", True),
        ("Вместо разработки отдельного фронтенда, система встраивается прямо в среду аналитика данных — Jupyter Notebook.", False),
        ("", False),
        ("Как это работает:", True),
        ("• Студент открывает TeachAI.ipynb и выполняет одну ячейку.", False),
        ("• Внизу ячейки отображается интерактивный интерфейс.", False),
        ("• Доступны кнопки: «Пройти тест», «Объясни подробнее», «Привести пример».", False),
        ("• Индикаторы загрузки (LoadingManager) обеспечивают плавный UX во время ожиданий ответа от LLM API.", False)
    ], Inches(1), Inches(1.5), Inches(11), Inches(5), font_size=18)

    # 11. Метрики (были 3 картинки)
    slide = slides[10]
    clear_picture_placeholders(slide)
    for shape in list(slide.shapes):
        if hasattr(shape, 'text'):
            if "Слайд с картинк" in shape.text: set_shape_text(shape, "Оценка качества системы", bold=True)
            elif "подпись" in shape.text.strip().lower() or "Пояснения" in shape.text:
                try: shape._element.getparent().remove(shape._element)
                except AttributeError: pass
                
    add_clean_textbox(slide, [
        ("95%", True),
        ("Валидность планов", False)
    ], Inches(1.0), Inches(2.5), Inches(3.0), Inches(2.0), font_size=40)
    add_clean_textbox(slide, [
        ("90%", True),
        ("Рабочий Python Код", False)
    ], Inches(5.0), Inches(2.5), Inches(3.0), Inches(2.0), font_size=40)
    add_clean_textbox(slide, [
        ("85%", True),
        ("Релевантность ответов", False)
    ], Inches(9.0), Inches(2.5), Inches(3.0), Inches(2.0), font_size=40)
    
    add_clean_textbox(slide, [
        ("Метрики получены путём ручного тестирования системы. Prompt engineering и строгая валидация JSON-схем обеспечивают стабильно высокое качество без дообучения.", False)
    ], Inches(1.0), Inches(5.5), Inches(11), Inches(1.5), font_size=16)

    # 12. Модульная структура
    slide = slides[11]
    clear_picture_placeholders(slide)
    for shape in list(slide.shapes):
        if hasattr(shape, 'text'):
            if "Слайд с картинк" in shape.text: set_shape_text(shape, "Модульная структура", bold=True)
            elif "Подпись" in shape.text:
                try: shape._element.getparent().remove(shape._element)
                except AttributeError: pass

    add_clean_textbox(slide, [("Ядро системы", True), ("• engine.py\n• main.py\n• logger.py", False)], Inches(1.5), Inches(2.5), Inches(2.5), Inches(2), font_size=18)
    add_clean_textbox(slide, [("Состояние", True), ("• state_manager.py\n• course_data.py\n• user_profile.py", False)], Inches(4.5), Inches(2.5), Inches(2.5), Inches(2), font_size=18)
    add_clean_textbox(slide, [("Лингвистика", True), ("• content_generator.py\n• 7 саб-генераторов\n• QA / Relevance", False)], Inches(7.5), Inches(2.5), Inches(2.5), Inches(2), font_size=18)
    add_clean_textbox(slide, [("Интерфейс", True), ("• interface.py\n• ipywidgets handlers\n• assessment.py", False)], Inches(10.5), Inches(2.5), Inches(2.5), Inches(2), font_size=18)

    # 13. Tech Stack
    slide = slides[12]
    for shape in slide.shapes:
        if hasattr(shape, 'text') and "большого объема" in shape.text.lower(): set_shape_text(shape, "Технологический стек", bold=True)
    add_clean_textbox(slide, [
        ("Основные технологии:", True),
        ("• Python 3.10+ — язык разработки всей системы", False),
        ("• OpenAI API (GPT-4 / GPT-3.5-turbo) — генерация контента", False),
        ("• Jupyter Notebook — среда выполнения", False),
        ("• ipywidgets 8.x — интерактивные виджеты для UI", False),
        ("• JSON — формат хранения данных и состояния", False),
        ("", False),
        ("Статистика проекта:", True),
        ("• Более 40 Python-модулей со строгим разделением ответственности", False),
        ("• 15,000+ строк кода", False),
        ("• Применение паттернов Facade, Strategy, Observer", False)
    ], Inches(1), Inches(1.5), Inches(11), Inches(5.5), font_size=16)

    # 14. QA
    slide = slides[13]
    for shape in slide.shapes:
        if hasattr(shape, 'text') and "большого объема" in shape.text.lower(): set_shape_text(shape, "Обеспечение качества (без дообучения)", bold=True)
    add_clean_textbox(slide, [
        ("1. Строгая инженерия промптов", True),
        ("   Контроль модели через четкие JSON-инструкции, эмоджи-маркеры и Few-Shot примеры.", False),
        ("2. Многоуровневая автокоррекция", True),
        ("   Система делает до 3 попыток генерации в случае сбоя парсинга JSON или нарушения структуры.", False),
        ("3. Проверка релевантности (Guardrails)", True),
        ("   Модуль relevance_checker анализирует вопросы студента, предотвращая отклонение LLM от темы.", False),
        ("4. Принудительная рандомизация", True),
        ("   Варианты ответов в тестах перемешиваются на стороне Python, так как LLM склонна к паттернам.", False)
    ], Inches(1), Inches(1.5), Inches(11), Inches(5.5), font_size=16)

    # 15. Justification
    slide = slides[14]
    for shape in slide.shapes:
        if hasattr(shape, 'text') and "большого объема" in shape.text.lower(): set_shape_text(shape, "Обоснование: Prompt Engineering vs Fine-Tuning", bold=True)
    add_clean_textbox(slide, [
        ("Техническая целесообразность:", True),
        ("Дообучение (Fine-tuning) LLM зафиксировало бы её 'стиль', но потребовало бы огромного датасета. Prompt Engineering позволяет модели гибко переключаться между темами.", False),
        ("", False),
        ("Экономическая эффективность:", True),
        ("Отсутствие затрат на разметку десятков тысяч примеров и использование мощных GPU. Инференс дешевле.", False),
        ("", False),
        ("Научное обоснование:", True),
        ("Исследования (Wei et al., 2022) доказывают, что методы Chain-of-Thought значительно улучшают логику LLM прямо 'из коробки', делая дообучение избыточным для общих образовательных задач.", False)
    ], Inches(1), Inches(1.5), Inches(11), Inches(5.5), font_size=16)

    # 16. Выводы
    slide = slides[15]
    for shape in slide.shapes:
        if hasattr(shape, 'text') and "Вывод" in shape.text: set_shape_text(shape, "Выводы", bold=True, font_size=28)
    
    add_clean_textbox(slide, [
        ("✅ Прототип готов:", True),
        ("Создана полнофункциональная образовательная система TeachAI, интегрируемая прямо в Jupyter Notebook.", False),
        ("", False),
        ("✅ Гипотеза о промптах подтверждена:", True),
        ("Грамотный Prompt Engineering полностью закрывает потребность в генерации валидного кода, качественных планов курсов и релевантности, обходя необходимость в Fine-Tuning.", False),
        ("", False),
        ("✅ Расширяемость:", True),
        ("Архитектура модулей-Фасадов позволяет легко добавить поддержку новых языков программирования или других LLM провайдеров.", False)
    ], Inches(1), Inches(1.5), Inches(11), Inches(5), font_size=18, bold_first=True)

    # 17. Заключение
    slide = slides[16]
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Заключение" in shape.text: set_shape_text(shape, "Заключительные слова и планы развития", bold=True, font_size=24)
            elif shape.text.startswith("Текст"):
                set_multiline_text(shape, [
                    "• Обертка проекта в полноценное веб-приложение (Flask/Django).",
                    "• Поддержка загрузки пользовательских документов (RAG - Retrieval Augmented Generation) для ответов по методичкам ВУЗа.",
                    "• Рекомендательная система на базе истории тестов студента.",
                    "",
                    "Спасибо за внимание!"
                ], font_size=18, color=RGBColor(0x33,0x33,0x33))
                shape.left = Inches(1)
                shape.width = Inches(11)

    prs.save(OUTPUT_PATH)
    print(f"✅ Presentation saved to: {OUTPUT_PATH}")

if __name__ == "__main__":
    create_presentation()
