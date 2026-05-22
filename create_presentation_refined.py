"""
Refined presentation generator for TeachAI project.

Creates a cleaner, more balanced deck from the existing template.
"""

# pyright: reportMissingImports=false

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
except ImportError as exc:
    raise SystemExit(
        "Missing dependency 'python-pptx'. "
        "Run this script with an environment where it is installed."
    ) from exc

TEMPLATE_PATH = "Шаблон презентации.pptx"
OUTPUT_PATH = "TeachAI_Stage4_refined.pptx"


def set_shape_text(shape, text, font_size=None, bold=None, color=None):
    """Replace text in shape while keeping most template formatting."""
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""

    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
    else:
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()

    run.text = text
    if font_size is not None:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def set_multiline(shape, lines, font_size=16, bold_first=False):
    """Write multiline bullet-like content into a placeholder."""
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]
        p._p.getparent().remove(p._p)

    first = tf.paragraphs[0]
    for run in first.runs:
        run.text = ""

    for i, line in enumerate(lines):
        p = first if i == 0 else tf.add_paragraph()
        run = p.runs[0] if p.runs else p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bool(bold_first and i == 0)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


def create_presentation():
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)

    # Slide 1
    for shape in slides[0].shapes:
        if hasattr(shape, "text"):
            if "Название" in shape.text:
                set_shape_text(
                    shape,
                    "TeachAI\nИнтеллектуальная система персонализированного обучения",
                    font_size=24,
                    bold=True,
                )
            elif "ФИО" in shape.text:
                set_shape_text(shape, "Игорь Пучкин", font_size=16)
            elif "Дата" in shape.text:
                set_shape_text(shape, "Апрель 2026", font_size=14)
            elif "Поток" in shape.text:
                set_shape_text(shape, "AI/ML — разработчик", font_size=14)

    # Slide 2
    for shape in slides[1].shapes:
        if hasattr(shape, "text"):
            if "Постановка" in shape.text:
                set_shape_text(shape, "1. Проблема и решение", bold=True)
            elif shape.text.startswith("Текст"):
                set_multiline(
                    shape,
                    [
                        "Проблема: создание персонализированного курса обычно требует много времени преподавателя.",
                        "",
                        "Решение: TeachAI автоматически генерирует план курса, уроки, примеры и тесты через OpenAI API.",
                        "",
                        "Ключевая идея: Prompt Engineering + валидация контента вместо fine-tuning.",
                    ],
                    font_size=14,
                )

    # Slide 3
    for shape in slides[2].shapes:
        if hasattr(shape, "text"):
            if "Цель" in shape.text:
                set_shape_text(shape, "2. Цель проекта", bold=True)
            elif shape.text.startswith("Текст"):
                set_multiline(
                    shape,
                    [
                        "Сделать масштабируемую систему обучения,",
                        "которая быстро создает качественный учебный контент",
                        "под параметры конкретного студента",
                        "в привычной среде Jupyter Notebook.",
                    ],
                    font_size=17,
                )

    # Slide 4
    for shape in slides[3].shapes:
        if hasattr(shape, "text"):
            if "Задачи" in shape.text:
                set_shape_text(shape, "3. Задачи проекта", bold=True)
            elif shape.text.strip().startswith("1"):
                set_multiline(
                    shape,
                    [
                        "1) Генерация учебного плана, уроков, примеров и тестов",
                        "2) Персонализация стиля и темпа обучения",
                        "3) Проверка качества и автокоррекция ответов LLM",
                        "4) Интерактивный интерфейс на ipywidgets",
                        "5) Сохранение прогресса и результатов",
                    ],
                    font_size=14,
                )

    # Slide 5
    for shape in slides[4].shapes:
        if hasattr(shape, "text"):
            if "база" in shape.text.lower():
                set_shape_text(shape, "4. Данные и модель", bold=True)
            elif "Объем" in shape.text or "информация" in shape.text.lower():
                set_multiline(
                    shape,
                    [
                        "Отдельный датасет не собирается: проект работает через предобученные GPT-модели.",
                        "",
                        "Источник генерации: OpenAI API (gpt-4 / gpt-3.5-turbo).",
                        "",
                        "Качество обеспечивается комбинацией:",
                        "• строгих промптов",
                        "• структурной и семантической валидации",
                        "• повторной генерации при ошибках",
                    ],
                    font_size=13,
                )

    # Slide 6
    for shape in slides[5].shapes:
        if hasattr(shape, "text") and "картинк" in shape.text.lower():
            set_shape_text(shape, "5. Архитектура системы", bold=True)

    # Slide 7
    for shape in slides[6].shapes:
        if hasattr(shape, "text"):
            if "картинк" in shape.text.lower():
                set_shape_text(shape, "6. Генерация контента (Facade)", bold=True)
            elif shape.text.startswith("Текст"):
                set_multiline(
                    shape,
                    [
                        "ContentGenerator координирует специализированные модули:",
                        "• CoursePlanGenerator",
                        "• LessonGenerator",
                        "• ExamplesGenerator",
                        "• AssessmentGenerator",
                        "• QAGenerator / ConceptsGenerator / RelevanceChecker",
                        "",
                        "Это упрощает масштабирование и поддержку проекта.",
                    ],
                    font_size=12,
                )

    # Slide 8
    for shape in slides[7].shapes:
        if hasattr(shape, "text"):
            if "картинк" in shape.text.lower():
                set_shape_text(shape, "7. Параметризация", bold=True)
            elif "Текст" in shape.text:
                set_multiline(
                    shape,
                    [
                        "Параметры LLM: model, temperature, max_tokens, response_format.",
                        "Параметры студента: стиль общения, время обучения, длительность урока.",
                        "Хранение состояния: data/state.json (прогресс, оценки, курс).",
                        "Поддерживаемые стили: friendly, formal, casual, brief.",
                    ],
                    font_size=12,
                )

    # Slide 9
    for shape in slides[8].shapes:
        if hasattr(shape, "text"):
            if "картинк" in shape.text.lower():
                set_shape_text(shape, "8. Жизненный цикл урока", bold=True)
            elif "Текст" in shape.text:
                set_multiline(
                    shape,
                    [
                        "1) Генерация урока и отображение в UI",
                        "2) Интерактив: объяснение, примеры, Q&A",
                        "3) Проверка релевантности вопросов к теме",
                        "4) Тест и контрольное задание",
                        "5) Сохранение результата и переход дальше",
                    ],
                    font_size=12,
                )

    # Slide 10
    for shape in slides[9].shapes:
        if hasattr(shape, "text") and "подпись" in shape.text.lower():
            set_shape_text(
                shape,
                "Точка входа: TeachAI.ipynb -> TeachAIEngine.start() -> интерактивный интерфейс",
                font_size=11,
            )

    # Slide 11
    for shape in slides[10].shapes:
        if hasattr(shape, "text"):
            if "Слайд с картинк" in shape.text:
                set_shape_text(shape, "9. Качество результата", bold=True)
            elif "Пояснения" in shape.text:
                set_multiline(
                    shape,
                    [
                        "Ключевые контрольные метрики качества контента:",
                        "• валидность структуры: ~95%",
                        "• корректность кода: ~90%",
                        "• релевантность ответов: ~85%",
                    ],
                    font_size=11,
                )

    # Slide 12
    for shape in slides[11].shapes:
        if hasattr(shape, "text"):
            if "Слайд с картинк" in shape.text:
                set_shape_text(shape, "10. Модульная структура", bold=True)
            elif "Подпись" in shape.text:
                # brief labels will be filled by template if needed
                pass

    # Slide 13
    for shape in slides[12].shapes:
        if hasattr(shape, "text") and "большого объема" in shape.text.lower():
            set_shape_text(shape, "11. Технологический стек", bold=True)
    slides[12].shapes.add_textbox(
        left=slides[12].shapes[0].left + 200000,
        top=slides[12].shapes[0].top + 1400000,
        width=9000000,
        height=3500000,
    ).text_frame.text = (
        "Python 3.10+, OpenAI API, Jupyter Notebook, ipywidgets, python-dotenv, httpx, JSON.\n"
        "Проект: 61 Python-файл, около 21k строк кода."
    )

    # Slide 14
    for shape in slides[13].shapes:
        if hasattr(shape, "text"):
            if "большого объема" in shape.text.lower():
                set_shape_text(shape, "12. Обеспечение качества", bold=True)
    slides[13].shapes.add_textbox(
        left=slides[13].shapes[0].left + 200000,
        top=slides[13].shapes[0].top + 1400000,
        width=9000000,
        height=3500000,
    ).text_frame.text = (
        "• Строгие промпты для каждого генератора\n"
        "• JSON/структурная валидация\n"
        "• Перегенерация при ошибках\n"
        "• RelevanceChecker для фильтрации вопросов\n"
        "• Логирование и отладочные ответы"
    )

    # Slide 15
    for shape in slides[14].shapes:
        if hasattr(shape, "text") and "большого объема" in shape.text.lower():
            set_shape_text(shape, "13. Почему без fine-tuning", bold=True)
    slides[14].shapes.add_textbox(
        left=slides[14].shapes[0].left + 200000,
        top=slides[14].shapes[0].top + 1400000,
        width=9000000,
        height=3500000,
    ).text_frame.text = (
        "Текущие задачи решаются prompt engineering подходом:\n"
        "• быстрое внедрение и низкая стоимость\n"
        "• гибкость по темам обучения\n"
        "• отсутствие необходимости в большом размеченном датасете\n"
        "• простое масштабирование и поддержка"
    )

    # Slide 16
    for shape in slides[15].shapes:
        if hasattr(shape, "text") and "Вывод" in shape.text:
            set_shape_text(shape, "14. Выводы", bold=True, font_size=28)

    # Slide 17
    for shape in slides[16].shapes:
        if hasattr(shape, "text"):
            if "Заключение" in shape.text:
                set_shape_text(shape, "15. Дальнейшее развитие", bold=True, font_size=24)
            elif shape.text.startswith("Текст"):
                set_multiline(
                    shape,
                    [
                        "• Веб-версия интерфейса (Flask/Django)",
                        "• A/B тестирование промптов",
                        "• Расширение каталога курсов",
                        "• Улучшенная аналитика обучения",
                        "",
                        "Спасибо за внимание!",
                    ],
                    font_size=14,
                )

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    create_presentation()
