"""
Скрипт для создания презентации Этапа 4 проекта TeachAI
на основе шаблона 'Шаблон презентации.pptx'
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

TEMPLATE_PATH = "Шаблон презентации.pptx"
OUTPUT_PATH = "TeachAI_Stage4.pptx"

# Paths to generated images
BRAIN_DIR = "/home/ip-59/.gemini/antigravity/brain/242bef71-a864-4600-8d99-4906f89eaa57"
IMG_ARCHITECTURE = os.path.join(BRAIN_DIR, [f for f in os.listdir(BRAIN_DIR) if f.startswith("architecture_diagram")][0])
IMG_LIFECYCLE = os.path.join(BRAIN_DIR, [f for f in os.listdir(BRAIN_DIR) if f.startswith("lesson_lifecycle")][0])
IMG_METRICS = os.path.join(BRAIN_DIR, [f for f in os.listdir(BRAIN_DIR) if f.startswith("quality_metrics")][0])
IMG_TECH = os.path.join(BRAIN_DIR, [f for f in os.listdir(BRAIN_DIR) if f.startswith("tech_stack")][0])
IMG_PROMPT = os.path.join(BRAIN_DIR, [f for f in os.listdir(BRAIN_DIR) if f.startswith("prompt_engineering")][0])


def clear_shape_text(shape):
    """Clear all text from a shape while preserving formatting of first run."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = ""


def set_shape_text(shape, text, font_size=None, bold=None, color=None, alignment=None):
    """Set text on a shape, preserving existing formatting where possible."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Clear existing paragraphs
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    
    # Set text on first paragraph's first run, or create one
    if tf.paragraphs and tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.text = text
    else:
        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        run = p.add_run()
        run.text = text
    
    if font_size is not None:
        run.font.size = Pt(font_size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    if alignment is not None:
        tf.paragraphs[0].alignment = alignment


def set_multiline_text(shape, lines, font_size=None, color=None, bold=None, line_spacing=None):
    """Set multiple lines of text on a shape."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    
    # Get formatting from first existing run
    existing_font_size = None
    existing_color = None
    existing_bold = None
    existing_font_name = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        r = tf.paragraphs[0].runs[0]
        existing_font_size = r.font.size
        try:
            existing_color = r.font.color.rgb if r.font.color and r.font.color.type is not None else None
        except AttributeError:
            existing_color = None
        existing_bold = r.font.bold
        existing_font_name = r.font.name
    
    # Clear existing text
    for i in range(len(tf.paragraphs) - 1, -1, -1):
        if i > 0:
            p = tf.paragraphs[i]
            p._p.getparent().remove(p._p)
    
    # Clear first paragraph
    first_para = tf.paragraphs[0]
    for run in first_para.runs:
        run.text = ""
    
    for i, line in enumerate(lines):
        if i == 0:
            p = first_para
            if p.runs:
                run = p.runs[0]
                run.text = line
            else:
                run = p.add_run()
                run.text = line
        else:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = line
        
        # Apply formatting
        fs = font_size if font_size is not None else existing_font_size
        if fs is not None:
            if isinstance(fs, (int, float)):
                run.font.size = Pt(fs)
            else:
                run.font.size = fs
        
        c = color if color is not None else existing_color
        if c is not None:
            run.font.color.rgb = c
        
        b = bold if bold is not None else existing_bold
        if b is not None:
            run.font.bold = b
        
        if existing_font_name:
            run.font.name = existing_font_name
        
        if line_spacing is not None:
            p.line_spacing = Pt(line_spacing)


def replace_placeholder_image(slide, shape_name, image_path):
    """Replace a placeholder image shape with a new image."""
    for shape in slide.shapes:
        if shape.name == shape_name:
            # Get position and size
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            # Remove old shape
            sp = shape._element
            sp.getparent().remove(sp)
            # Add new image
            slide.shapes.add_picture(image_path, left, top, width, height)
            return True
    return False


def add_image_to_slide(slide, image_path, left, top, width, height=None):
    """Add an image to a slide at specific position."""
    if height:
        slide.shapes.add_picture(image_path, left, top, width, height)
    else:
        slide.shapes.add_picture(image_path, left, top, width=width)


def find_shape_by_text(slide, text_contains):
    """Find a shape containing specific text."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            if text_contains.lower() in shape.text.lower():
                return shape
    return None


def find_shape_by_name(slide, name):
    """Find a shape by name."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def create_presentation():
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)
    
    print(f"Template has {len(slides)} slides")
    
    # ================================================================
    # SLIDE 1: Титульный слайд
    # ================================================================
    slide = slides[0]
    print("Processing Slide 1: Title")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Название" in shape.text:
                set_shape_text(shape, "TeachAI — Интеллектуальная система\nперсонализированного обучения", font_size=24, bold=True)
            elif "ФИО" in shape.text:
                set_shape_text(shape, "Игорь Пучкин", font_size=16)
            elif "Дата" in shape.text:
                set_shape_text(shape, "Февраль 2026", font_size=14)
            elif "Поток" in shape.text:
                set_multiline_text(shape, ["\n", "AI/ML — разработчик"], font_size=14)
    
    # ================================================================
    # SLIDE 2: Постановка задачи
    # ================================================================
    slide = slides[1]
    print("Processing Slide 2: Problem Statement")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Постановка" in shape.text:
                set_shape_text(shape, "1. Тема и описание задачи", bold=True)
            elif shape.text.startswith("Текст"):
                set_multiline_text(shape, [
                    "Проблема:",
                    "Отсутствие доступных инструментов для персонализированного",
                    "обучения программированию с адаптацией под студента.",
                    "",
                    "Решение: TeachAI — интеллектуальная образовательная система,",
                    "использующая GPT модель через OpenAI API для автоматической",
                    "генерации персонализированных учебных материалов.",
                    "",
                    "Тип задачи: AI-система прикладного уровня, использующая",
                    "inference предобученной LLM для генерации образовательного контента.",
                    "",
                    "Подход: Prompt Engineering вместо fine-tuning.",
                ], font_size=12)
    
    # ================================================================
    # SLIDE 3: Цель
    # ================================================================
    slide = slides[2]
    print("Processing Slide 3: Goal")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Цель" in shape.text:
                set_shape_text(shape, "Цель проекта", bold=True)
            elif shape.text.startswith("Текст"):
                set_multiline_text(shape, [
                    "Автоматизация создания персонализированного",
                    "образовательного контента с использованием",
                    "современных языковых моделей (LLM)",
                    "для обеспечения эффективного",
                    "и гибкого обучения в интерактивной среде",
                    "Jupyter Notebook.",
                ], font_size=16)
    
    # ================================================================
    # SLIDE 4: Задачи
    # ================================================================
    slide = slides[3]
    print("Processing Slide 4: Tasks")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Задачи" in shape.text:
                set_shape_text(shape, "Задачи проекта", bold=True)
            elif shape.text.strip().startswith("1"):
                set_multiline_text(shape, [
                    "1. Генерация учебных планов, уроков, примеров и тестов",
                    "    через OpenAI GPT API",
                    "",
                    "2. Персонализация контента под стиль общения",
                    "    и параметры обучения студента",
                    "",
                    "3. Многоуровневая валидация и автокоррекция",
                    "    сгенерированного контента",
                    "",
                    "4. Интерактивный интерфейс на ipywidgets",
                    "    с навигацией, тестированием и Q&A",
                    "",
                    "5. Отслеживание прогресса обучения",
                    "    и сохранение состояния",
                ], font_size=12)
    
    # ================================================================
    # SLIDE 5: Обучающая база
    # ================================================================
    slide = slides[4]
    print("Processing Slide 5: Training Data")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "база" in shape.text.lower():
                set_shape_text(shape, "2. База данных / Данные", bold=True)
            elif "Объем" in shape.text or "информация" in shape.text.lower():
                set_multiline_text(shape, [
                    "Сбор базы данных не требуется",
                    "Проект использует готовую предобученную GPT-модель через API",
                    "",
                    "Источник данных:",
                    "OpenAI GPT-4 API — генерация контента в реальном времени",
                    "",
                    "Вместо датасета:",
                    "Prompt Engineering — управление качеством через промпты",
                    "",
                    "Обоснование:",
                    "• Нет обучения/дообучения модели",
                    "• Нет задач классификации, регрессии, detection",
                    "• Контент генерируется динамически по запросу",
                    "• Подход обоснован на Этапах 2 и 3",
                ], font_size=11)
    
    # ================================================================
    # SLIDE 6: Архитектура (картинка на весь слайд)
    # ================================================================
    slide = slides[5]
    print("Processing Slide 6: Architecture Overview")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and "картинк" in shape.text.lower():
            set_shape_text(shape, "4. Архитектура системы TeachAI", bold=True)
    
    # Add architecture diagram image
    add_image_to_slide(slide, IMG_ARCHITECTURE,
                       Inches(1.5), Inches(1.2), Inches(7), Inches(5.5))
    
    # ================================================================
    # SLIDE 7: Архитектура генерации контента
    # ================================================================
    slide = slides[6]
    print("Processing Slide 7: Content Generation Architecture")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "картинк" in shape.text.lower():
                set_shape_text(shape, "Prompt Engineering — методология", bold=True)
            elif shape.text.startswith("Текст"):
                set_multiline_text(shape, [
                    "ContentGenerator (фасад):",
                    "",
                    "• CoursePlanGenerator",
                    "• LessonGenerator",
                    "• ExamplesGenerator",
                    "• AssessmentGenerator",
                    "• QAGenerator",
                    "• ConceptsGenerator",
                    "• RelevanceChecker",
                    "",
                    "Каждый генератор использует",
                    "специализированные промпты",
                    "для максимального качества",
                ], font_size=10)
    
    # Find and replace image placeholder
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top, w, h = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(IMG_PROMPT, left, top, w, h)
            break
    
    # ================================================================
    # SLIDE 8: Параметризация данных
    # ================================================================
    slide = slides[7]
    print("Processing Slide 8: Data Parameterization")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "картинк" in shape.text.lower():
                set_shape_text(shape, "3. Параметризация данных", bold=True)
            elif "Текст" in shape.text:
                set_multiline_text(shape, [
                    "Параметры API запросов:",
                    "• model: gpt-4 / gpt-3.5-turbo",
                    "• temperature: 0.7 (генерация) / 0.2 (проверка)",
                    "• max_tokens: 3500",
                    "• response_format: json_object",
                    "",
                    "Параметры профиля студента:",
                    "• Стиль общения (formal/friendly/casual/brief)",
                    "• Время обучения (часы)",
                    "• Длительность урока (минуты)",
                    "",
                    "Структура данных:",
                    "• JSON-схемы для планов курсов",
                    "• state.json — персистентное состояние",
                    "• Кэширование контента уроков",
                ], font_size=10)
    
    # Replace image with lifecycle diagram  
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top, w, h = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(IMG_LIFECYCLE, left, top, w, h)
            break
    
    # ================================================================
    # SLIDE 9: Жизненный цикл (слайд с картинкой и текстом)
    # ================================================================
    slide = slides[8]
    print("Processing Slide 9: Lesson Lifecycle")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "картинк" in shape.text.lower():
                set_shape_text(shape, "5. Жизненный цикл урока", bold=True)
            elif "Текст" in shape.text:
                set_multiline_text(shape, [
                    "Каждый урок проходит цикл:",
                    "",
                    "1. Загрузка → Генерация контента через GPT",
                    "2. Отображение материала студенту",
                    "3. Интерактивное взаимодействие:",
                    "   • Объяснение подробнее",
                    "   • Примеры кода",
                    "   • Вопросы по теме",
                    "4. Тестирование (5 вопросов)",
                    "5. Контрольное задание",
                    "6. Переход к следующему уроку",
                ], font_size=11)
    
    # ================================================================
    # SLIDE 10: Скриншот (фото слайд)
    # ================================================================
    slide = slides[9]
    print("Processing Slide 10: Screenshot/Interface")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text') and "подпись" in shape.text.lower():
            set_shape_text(shape, "Точка входа: TeachAI.ipynb → engine.py → start()", font_size=11)
    
    # Replace the main image with tech stack
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top, w, h = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(IMG_TECH, left, top, w, h)
            break
    
    # ================================================================
    # SLIDE 11: 3 картинки (метрики)
    # ================================================================
    slide = slides[10]
    print("Processing Slide 11: Quality Metrics (3 images)")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Слайд с картинк" in shape.text:
                set_shape_text(shape, "5. Графическое подтверждение — метрики качества", bold=True)
            elif "Пояснения" in shape.text:
                set_multiline_text(shape, [
                    "Метрики получены путём тестирования системы без дообучения модели.",
                    "Prompt engineering + валидация обеспечивают стабильное качество.",
                ], font_size=11)
            elif "подпись" == shape.text.strip().lower():
                # There are 3 caption fields
                pass
    
    # Replace the 3 placeholder images with the metrics image
    replaced = 0
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            left, top, w, h = shape.left, shape.top, shape.width, shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            if replaced == 0:
                # First image slot — use the full metrics image spanning all 3 slots
                slide.shapes.add_picture(IMG_METRICS, left, top, Inches(10), h)
            replaced += 1
    
    # Update captions
    caption_idx = 0
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip().lower() == "подпись":
            captions = ["Валидность планов: ~95%", "Корректность кода: ~90%", "Релевантность: ~85%"]
            if caption_idx < len(captions):
                set_shape_text(shape, captions[caption_idx], font_size=10)
                caption_idx += 1
    
    # ================================================================
    # SLIDE 12: 4 картинки (модульная структура)
    # ================================================================
    slide = slides[11]
    print("Processing Slide 12: Module Structure")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Слайд с картинк" in shape.text:
                set_shape_text(shape, "6. Ноутбук / модульная структура проекта", bold=True)
    
    # Update caption texts
    captions_data = [
        "Ядро: engine.py,\nconfig.py, logger.py",
        "Состояние: state_manager,\nuser_profile, progress",
        "Генерация: content_generator,\n7 специализированных модулей",
        "UI: interface.py,\nlesson_*, assessment_*"
    ]
    caption_idx = 0
    for shape in slide.shapes:
        if hasattr(shape, 'text') and "Подпись" in shape.text:
            if caption_idx < len(captions_data):
                set_multiline_text(shape, captions_data[caption_idx].split("\n"), font_size=9)
                caption_idx += 1
    
    # ================================================================
    # SLIDE 13: Большой объём текста — Технологический стек
    # ================================================================
    slide = slides[12]
    print("Processing Slide 13: Tech Stack")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "большого объема" in shape.text.lower():
                set_shape_text(shape, "Технологический стек", bold=True)
    
    # Add a text box with the tech stack info
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(10.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    tech_lines = [
        ("Основные технологии:", True, 16),
        ("", False, 12),
        ("• Python 3.10+ — язык разработки всей системы", False, 13),
        ("• OpenAI API (GPT-4 / GPT-3.5-turbo) — генерация контента", False, 13),
        ("• Jupyter Notebook — среда выполнения и интерфейс", False, 13),
        ("• ipywidgets 8.x — интерактивные виджеты для UI", False, 13),
        ("• python-dotenv — управление конфигурацией", False, 13),
        ("• httpx — HTTP-клиент для работы с прокси", False, 13),
        ("• JSON — формат хранения данных и состояния", False, 13),
        ("", False, 12),
        ("Статистика проекта:", True, 16),
        ("", False, 12),
        ("• 40+ Python-модулей специализированной функциональности", False, 13),
        ("• 15,000+ строк кода", False, 13),
        ("• 7 генераторов контента (фасадный паттерн)", False, 13),
        ("• Точка входа: TeachAI.ipynb → engine.py → start()", False, 13),
    ]
    
    for i, (text, is_bold, size) in enumerate(tech_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    # ================================================================
    # SLIDE 14: Большой объём текста — Методы обеспечения качества
    # ================================================================
    slide = slides[13]
    print("Processing Slide 14: Quality Assurance Methods")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "большого объема" in shape.text.lower():
                set_shape_text(shape, "Методы обеспечения качества без дообучения", bold=True)
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(10.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    quality_lines = [
        ("1. Строгая инженерия промптов", True, 14),
        ("   Детальные инструкции с эмоджи-маркерами, многократное повторение", False, 12),
        ("   требований, примеры правильного/неправильного форматирования", False, 12),
        ("", False, 10),
        ("2. Многоуровневая валидация", True, 14),
        ("   • Структурная: проверка JSON-схем, обязательных полей", False, 12),
        ("   • Семантическая: проверка релевантности содержания теме", False, 12),
        ("   • Формат-валидация: корректность кода, markdown-разметки", False, 12),
        ("", False, 10),
        ("3. Автоматическая коррекция и перегенерация", True, 14),
        ("   При обнаружении ошибок — исправление структуры данных,", False, 12),
        ("   перегенерация с усиленными промптами, fallback-стратегии", False, 12),
        ("", False, 10),
        ("4. Проверка релевантности в реальном времени", True, 14),
        ("   Модуль relevance_checker.py анализирует вопросы студентов", False, 12),
        ("   на соответствие теме урока через API (temperature=0.2)", False, 12),
        ("", False, 10),
        ("5. Рандомизация тестов", True, 14),
        ("   Принудительная рандомизация вариантов ответов, т.к. API", False, 12),
        ("   не всегда следует инструкциям по случайному порядку", False, 12),
    ]
    
    for i, (text, is_bold, size) in enumerate(quality_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    # ================================================================
    # SLIDE 15: Большой объём текста — Обоснование подхода
    # ================================================================
    slide = slides[14]
    print("Processing Slide 15: Approach Justification")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "большого объема" in shape.text.lower():
                set_shape_text(shape, "Обоснование подхода: Prompt Engineering vs Fine-tuning", bold=True)
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(10.5), Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    justification_lines = [
        ("Почему НЕ используется дообучение модели:", True, 14),
        ("", False, 10),
        ("Технические причины:", True, 13),
        ("• Гибкость: система работает с любой темой без переобучения", False, 12),
        ("• Актуальность: автоматическое улучшение при обновлении модели OpenAI", False, 12),
        ("• Малый объём данных: недостаточно примеров для fine-tuning", False, 12),
        ("", False, 10),
        ("Экономические причины:", True, 13),
        ("• Нет затрат на GPU для обучения", False, 12),
        ("• Нет затрат на создание и разметку датасетов", False, 12),
        ("", False, 10),
        ("Научное обоснование:", True, 13),
        ("• OpenAI рекомендует начинать с промптов, fine-tuning — при необходимости", False, 12),
        ("• Wei et al., 2022: Chain-of-Thought Prompting улучшает качество без дообучения", False, 12),
        ("• Wang et al., 2022: Self-Consistency — валидация повышает надёжность LLM", False, 12),
        ("", False, 10),
        ("Fine-tuning оправдан только при: 10,000+ примеров, узкая специализация,", False, 12),
        ("критичная стабильность формата, тысячи запросов в день.", False, 12),
        ("Ни одно из условий не выполняется в данном проекте.", True, 12),
    ]
    
    for i, (text, is_bold, size) in enumerate(justification_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    
    # ================================================================
    # SLIDE 16: Выводы
    # ================================================================
    slide = slides[15]
    print("Processing Slide 16: Conclusions")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Вывод" in shape.text:
                set_shape_text(shape, "7. Выводы", bold=True, font_size=28)
    
    # Add conclusions text box
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(6), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    conclusions = [
        ("Результаты:", True, 15),
        ("", False, 10),
        ("✅ Создан работающий прототип интеллектуальной", False, 13),
        ("    образовательной системы на базе GPT", False, 13),
        ("", False, 10),
        ("✅ Prompt Engineering обеспечивает качество:", False, 13),
        ("    • Валидность планов: ~95%", False, 13),
        ("    • Корректность кода: ~90%", False, 13),
        ("    • Релевантность контента: ~85%", False, 13),
        ("", False, 10),
        ("✅ Интерактивный интерфейс с полным циклом", False, 13),
        ("    обучения в Jupyter Notebook", False, 13),
        ("", False, 10),
        ("✅ Дообучение модели не требуется —", False, 13),
        ("    промпт-инженерия решает все задачи", False, 13),
    ]
    
    for i, (text, is_bold, size) in enumerate(conclusions):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    
    # ================================================================
    # SLIDE 17: Заключение / План дальнейшей работы
    # ================================================================
    slide = slides[16]
    print("Processing Slide 17: Future Work / Conclusion")
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            if "Заключение" in shape.text:
                set_shape_text(shape, "8. План дальнейшей работы", bold=True, font_size=24)
            elif shape.text.startswith("Текст"):
                set_multiline_text(shape, [
                    "• Веб-интерфейс (Flask/Django) для доступа без Jupyter",
                    "• Мультимодальность: диаграммы, визуализации в уроках",
                    "• Система рекомендаций на основе прогресса студента",
                    "• A/B тестирование промптов для оптимизации качества",
                    "• Расширение каталога курсов",
                    "",
                    "Спасибо за внимание!",
                ], font_size=14)
    
    # ================================================================
    # Save the presentation
    # ================================================================
    prs.save(OUTPUT_PATH)
    print(f"\n✅ Presentation saved to: {OUTPUT_PATH}")
    print(f"   Total slides: {len(slides)}")


if __name__ == "__main__":
    create_presentation()
